"""Build the 'Four Homepage Archetypes' deck.

Source: _research/landing-pages-2026-04.md
Output: _research/decks/angel-syndicate-archetypes.pptx

Brand palette: cream / moss / sage / rust.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Brand
# ---------------------------------------------------------------------------
CREAM = RGBColor(0xFB, 0xF8, 0xF2)
MOSS  = RGBColor(0x2F, 0x3E, 0x34)
SAGE  = RGBColor(0x7E, 0x8E, 0x84)
RUST  = RGBColor(0xD1, 0x6B, 0x4F)
PLUM  = RGBColor(0x8B, 0x5E, 0x83)
SAGE_PALE = RGBColor(0xE3, 0xE7, 0xE4)   # very light sage tint for cards
MOSS_DEEP = RGBColor(0x1F, 0x2A, 0x23)

HEADER_FONT = "Cormorant Garamond"
BODY_FONT   = "Inter"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    return tb, tf


def style_run(run, *, font=BODY_FONT, size=14, color=MOSS, bold=False,
              italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def write_para(tf, text, *, first=False, font=BODY_FONT, size=14, color=MOSS,
               bold=False, italic=False, align=PP_ALIGN.LEFT, space_after=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    style_run(run, font=font, size=size, color=color, bold=bold, italic=italic)
    return p


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_amp_watermark(slide, x, y, size=10, color=None):
    """Decorative oversized ampersand. Brand watermark."""
    if color is None:
        color = SAGE_PALE
    tb, tf = add_text_box(slide, x, y, size, size)
    write_para(
        tf, "&",
        first=True, font=HEADER_FONT, size=520, color=color, italic=False,
    )


def add_footer(slide, label):
    tb, tf = add_text_box(slide, 0.6, 7.05, 12.1, 0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "Silver & Salt Capital"
    style_run(r1, size=9, color=SAGE, italic=True)
    r2 = p.add_run()
    r2.text = "   ·   "
    style_run(r2, size=9, color=SAGE)
    r3 = p.add_run()
    r3.text = label
    style_run(r3, size=9, color=SAGE)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, MOSS_DEEP)

    # giant cream ampersand watermark, lower right
    tb, tf = add_text_box(slide, 7.5, 0.5, 6, 7)
    write_para(
        tf, "&", first=True, font=HEADER_FONT, size=560,
        color=RGBColor(0x3A, 0x4B, 0x40), italic=False, align=PP_ALIGN.LEFT,
    )

    # eyebrow
    tb, tf = add_text_box(slide, 0.7, 1.4, 7.5, 0.4)
    write_para(
        tf, "LANDSCAPE STUDY  ·  APRIL 2026", first=True,
        font=BODY_FONT, size=11, color=SAGE, bold=True,
    )

    # title
    tb, tf = add_text_box(slide, 0.7, 2.0, 9.5, 2.6)
    write_para(
        tf, "The four homepage archetypes", first=True,
        font=HEADER_FONT, size=66, color=CREAM, italic=True,
    )

    # subtitle
    tb, tf = add_text_box(slide, 0.7, 4.5, 10, 1.2)
    write_para(
        tf,
        "How 53 women-focused angel networks and syndicates introduce themselves online,",
        first=True, font=HEADER_FONT, size=22, color=CREAM,
    )
    write_para(
        tf,
        "and what the patterns reveal about what is missing.",
        font=HEADER_FONT, size=22, color=CREAM, italic=True,
    )

    # rule + caption
    add_rect(slide, 0.7, 6.4, 1.0, 0.03, fill=RUST)
    tb, tf = add_text_box(slide, 0.7, 6.55, 10, 0.4)
    write_para(
        tf,
        "Method: 49 of 53 homepages parsed for title, headings, navigation, CTAs, lead copy, and tech-stack signals.",
        first=True, font=BODY_FONT, size=11, color=SAGE, italic=True,
    )
    return slide


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    # title
    tb, tf = add_text_box(slide, 0.7, 0.55, 12, 0.7)
    write_para(
        tf, "Four patterns dominate the field.", first=True,
        font=HEADER_FONT, size=38, color=MOSS, italic=True,
    )

    # rule
    add_rect(slide, 0.7, 1.3, 1.0, 0.03, fill=RUST)

    tb, tf = add_text_box(slide, 0.7, 1.45, 12, 0.45)
    write_para(
        tf,
        "Across 40 working homepages, every site fell into one of four shapes. The mix is lopsided.",
        first=True, font=BODY_FONT, size=14, color=SAGE, italic=True,
    )

    # four columns
    cols = [
        ("01", "Two-door homepage", "~22 sites", "Founders and investors share equal billing in the hero."),
        ("02", "Investor-only club", "~6 sites",  "Page is pitched at prospective members, founders are absent."),
        ("03", "Platform / multi-product", "~5 sites", "Less a network, more a content and tools brand."),
        ("04", "Personal / portfolio", "2 sites",  "Lower-key, individual-investor presence."),
    ]
    col_w = 2.95
    gap = 0.15
    start_x = 0.7
    top = 2.05
    h = 4.7

    for i, (num, name, count, gist) in enumerate(cols):
        x = start_x + i * (col_w + gap)
        # card body
        add_rect(slide, x, top, col_w, h, fill=RGBColor(0xFF, 0xFD, 0xF7))
        # accent bar
        accent = RUST if i == 0 else SAGE
        add_rect(slide, x, top, col_w, 0.06, fill=accent)

        # number
        tb, tf = add_text_box(slide, x + 0.25, top + 0.2, col_w - 0.5, 0.7)
        write_para(
            tf, num, first=True, font=HEADER_FONT, size=44, color=accent,
            italic=True,
        )

        # name
        tb, tf = add_text_box(slide, x + 0.25, top + 1.1, col_w - 0.5, 1.1)
        write_para(
            tf, name, first=True, font=HEADER_FONT, size=24, color=MOSS,
            italic=True,
        )

        # count chip
        tb, tf = add_text_box(slide, x + 0.25, top + 2.25, col_w - 0.5, 0.35)
        write_para(
            tf, count.upper(), first=True, font=BODY_FONT, size=10,
            color=SAGE, bold=True,
        )

        # gist
        tb, tf = add_text_box(slide, x + 0.25, top + 2.75, col_w - 0.5, 1.7)
        write_para(
            tf, gist, first=True, font=BODY_FONT, size=13, color=MOSS,
        )

    add_footer(slide, "The four homepage archetypes")
    return slide


def archetype_slide(prs, *, num, name, count, summary, examples, why_present,
                    why_chosen):
    """A detailed two-column slide for one archetype."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    # eyebrow
    tb, tf = add_text_box(slide, 0.7, 0.55, 5, 0.4)
    write_para(
        tf, f"ARCHETYPE {num}  ·  {count.upper()}", first=True,
        font=BODY_FONT, size=11, color=RUST, bold=True,
    )

    # title
    tb, tf = add_text_box(slide, 0.7, 0.95, 12, 1.0)
    write_para(
        tf, name, first=True, font=HEADER_FONT, size=42, color=MOSS,
        italic=True,
    )

    # rule
    add_rect(slide, 0.7, 2.0, 1.0, 0.03, fill=RUST)

    # summary
    tb, tf = add_text_box(slide, 0.7, 2.15, 11.9, 0.7)
    write_para(
        tf, summary, first=True, font=HEADER_FONT, size=20, color=MOSS,
    )

    # ---- left column: WHY this exists ----
    left_x = 0.7
    col_top = 3.0
    col_w = 5.7
    col_h = 3.2

    add_rect(slide, left_x, col_top, col_w, col_h,
             fill=RGBColor(0xFF, 0xFD, 0xF7))
    add_rect(slide, left_x, col_top, 0.06, col_h, fill=RUST)

    tb, tf = add_text_box(slide, left_x + 0.3, col_top + 0.25, col_w - 0.6, 0.5)
    write_para(
        tf, "WHY IT EXISTS", first=True, font=BODY_FONT, size=10,
        color=RUST, bold=True,
    )
    tb, tf = add_text_box(slide, left_x + 0.3, col_top + 0.65, col_w - 0.6,
                          col_h - 0.85)
    write_para(
        tf, why_present, first=True, font=BODY_FONT, size=14, color=MOSS,
    )

    # ---- right column: WHO PICKS IT ----
    right_x = left_x + col_w + 0.3
    add_rect(slide, right_x, col_top, col_w, col_h,
             fill=RGBColor(0xFF, 0xFD, 0xF7))
    add_rect(slide, right_x, col_top, 0.06, col_h, fill=SAGE)

    tb, tf = add_text_box(slide, right_x + 0.3, col_top + 0.25, col_w - 0.6,
                          0.5)
    write_para(
        tf, "WHO PICKS IT", first=True, font=BODY_FONT, size=10,
        color=SAGE, bold=True,
    )
    tb, tf = add_text_box(slide, right_x + 0.3, col_top + 0.65, col_w - 0.6,
                          col_h - 0.85)
    write_para(
        tf, why_chosen, first=True, font=BODY_FONT, size=14, color=MOSS,
    )

    # ---- examples strip ----
    ex_top = col_top + col_h + 0.2
    tb, tf = add_text_box(slide, 0.7, ex_top, 2.0, 0.3)
    write_para(
        tf, "EXAMPLES", first=True, font=BODY_FONT, size=10, color=SAGE,
        bold=True,
    )
    tb, tf = add_text_box(slide, 0.7, ex_top + 0.3, 11.9, 0.5)
    write_para(
        tf, examples, first=True, font=BODY_FONT, size=12, color=MOSS,
        italic=True,
    )

    add_footer(slide, f"Archetype {num} of 4")
    return slide


def slide_takeaway(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, MOSS_DEEP)

    # eyebrow
    tb, tf = add_text_box(slide, 0.7, 0.6, 9, 0.4)
    write_para(
        tf, "WHAT THIS MEANS FOR SILVER & SALT CAPITAL", first=True,
        font=BODY_FONT, size=11, color=RUST, bold=True,
    )

    # title
    tb, tf = add_text_box(slide, 0.7, 1.05, 12, 1.4)
    write_para(
        tf, "The field has habits.", first=True, font=HEADER_FONT, size=44,
        color=CREAM, italic=True,
    )
    write_para(
        tf, "We have room to be specific.", font=HEADER_FONT, size=44,
        color=CREAM, italic=True,
    )

    # gaps grid (2 x 3)
    gaps = [
        ("Cost transparency",
         "Almost no site publishes dues, LP minimums, or per-deal commitments on the homepage."),
        ("Calendar booking on apply",
         "Most rely on long forms or email. Inline calendar bookings are rare."),
        ("Editorial output",
         "Aside from The Helm and Hustle Fund, very few sites publish original research."),
        ("Plain-language voice",
         "Most copy reads corporate or legal. The few warm voices stand out immediately."),
        ("Mobile fluency",
         "Many sites visibly have not been touched for mobile in years."),
        ("Visible LP economics",
         "Industry-standard 2 and 20 is assumed everywhere, stated nowhere."),
    ]
    cols = 3
    rows = 2
    grid_top = 3.0
    grid_left = 0.7
    cell_w = 4.0
    cell_h = 1.85
    gap = 0.2

    for i, (head, body) in enumerate(gaps):
        r = i // cols
        c = i % cols
        x = grid_left + c * (cell_w + gap)
        y = grid_top + r * (cell_h + gap)

        add_rect(slide, x, y, cell_w, cell_h, fill=RGBColor(0x3A, 0x4B, 0x40))
        add_rect(slide, x, y, 0.05, cell_h, fill=RUST)

        tb, tf = add_text_box(slide, x + 0.25, y + 0.18, cell_w - 0.5, 0.45)
        write_para(
            tf, head, first=True, font=HEADER_FONT, size=18, color=CREAM,
            italic=True,
        )
        tb, tf = add_text_box(slide, x + 0.25, y + 0.7, cell_w - 0.5,
                              cell_h - 0.85)
        write_para(
            tf, body, first=True, font=BODY_FONT, size=12,
            color=RGBColor(0xCF, 0xD2, 0xCC),
        )

    add_footer(slide, "Field gaps as differentiators")
    return slide


def slide_standouts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    tb, tf = add_text_box(slide, 0.7, 0.55, 12, 0.7)
    write_para(
        tf, "Worth studying.", first=True, font=HEADER_FONT, size=38,
        color=MOSS, italic=True,
    )
    add_rect(slide, 0.7, 1.3, 1.0, 0.03, fill=RUST)
    tb, tf = add_text_box(slide, 0.7, 1.45, 12, 0.45)
    write_para(
        tf,
        "Six pages do something distinctive that translates to our own join.html.",
        first=True, font=BODY_FONT, size=14, color=SAGE, italic=True,
    )

    rows = [
        ("Golden Seeds",
         "Cleanest two-door split: founder and member sub-navigation each have their own ladder."),
        ("Brydge Club",
         "Tightest investor-only pitch in the set. Cohort-based application reads as scarcity-driven."),
        ("Hustle Fund / Angel Squad",
         "Voice and content depth: newsletters, books, events, blog. Free resources as a funnel."),
        ("The Helm",
         "Combines fund, syndicate, and editorial brand. Single CTA ladder for founders and LPs."),
        ("Frontier Angels",
         "Strong regional angle. The 'we do not invest in' section saves both sides time."),
        ("Tidal River Fund",
         "Crisp tagline, honest about scale (a $1M fund), open geography despite the regional name."),
    ]
    rows_top = 2.05
    row_h = 0.78
    for i, (label, body) in enumerate(rows):
        y = rows_top + i * row_h
        # left badge
        add_rect(slide, 0.7, y + 0.08, 0.18, 0.5, fill=RUST if i % 2 == 0
                 else SAGE)
        # label
        tb, tf = add_text_box(slide, 1.0, y + 0.05, 3.4, 0.6,
                              anchor=MSO_ANCHOR.MIDDLE)
        write_para(
            tf, label, first=True, font=HEADER_FONT, size=22, color=MOSS,
            italic=True,
        )
        # body
        tb, tf = add_text_box(slide, 4.55, y + 0.05, 8.2, 0.65,
                              anchor=MSO_ANCHOR.MIDDLE)
        write_para(
            tf, body, first=True, font=BODY_FONT, size=13, color=MOSS,
        )

    add_footer(slide, "Six pages worth a closer read")
    return slide


def slide_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    # large quote
    tb, tf = add_text_box(slide, 1.2, 1.8, 11, 3.2)
    write_para(
        tf, "Most sites in the field are built around a fork in the road.",
        first=True, font=HEADER_FONT, size=34, color=MOSS, italic=True,
    )
    write_para(
        tf, "We are building one path, with the numbers visible.",
        font=HEADER_FONT, size=34, color=MOSS, italic=True,
    )

    add_rect(slide, 1.2, 5.3, 1.0, 0.04, fill=RUST)
    tb, tf = add_text_box(slide, 1.2, 5.5, 11, 0.4)
    write_para(
        tf, "Source: _research/landing-pages-2026-04.md", first=True,
        font=BODY_FONT, size=11, color=SAGE, italic=True,
    )

    # corner ampersand
    tb, tf = add_text_box(slide, 9.5, 4.5, 4, 4)
    write_para(
        tf, "&", first=True, font=HEADER_FONT, size=420,
        color=RGBColor(0xEC, 0xE7, 0xDB),
    )
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_overview(prs)

    archetype_slide(
        prs,
        num="01",
        name="The two-door homepage",
        count="~22 sites · most common",
        summary="The hero splits into 'For Founders' and 'For Investors / Members' with equal billing.",
        examples="Golden Seeds, Ark Angel Alliance, Tidal River Fund, TCA Venture Group, Launchpad, Maine Angels, Citrine Angels, NY Angels, Frontier Angels, Nebraska Angels, JumpFund, BEAM, Belle Capital.",
        why_present=(
            "These groups serve two audiences with one budget: founders seeking checks "
            "and accredited members seeking deal flow. A split hero is the cheapest way "
            "to acknowledge both without choosing. Section patterns repeat almost word "
            "for word: Mission, Investment Criteria, Portfolio, Apply for Funding, "
            "Become a Member, Events, News."
        ),
        why_chosen=(
            "Established networks that have run for a decade or more, regional angel "
            "groups with steady deal flow, and groups whose member count is the product "
            "they sell. A two-door homepage signals 'we are a real institution and you "
            "are welcome on either side,' which tracks for a board-led nonprofit or "
            "professionally managed network."
        ),
    )

    archetype_slide(
        prs,
        num="02",
        name="The investor-only club",
        count="~6 sites",
        summary="The founders application is hidden or absent. The page speaks only to prospective members.",
        examples="Brydge Club, The Beam Network, Plum Alley Ventures, Impact Invest Her, Broadway Angels, OSEA.",
        why_present=(
            "The product these groups actually sell is membership: education, vetted "
            "deal flow, and peer community. Putting a founder application alongside it "
            "muddies the offer. By cutting the founder lane, the page becomes a sales "
            "page for one thing, with copy that reads more like a private club "
            "application than a public investment portal."
        ),
        why_chosen=(
            "Smaller, newer collectives where the membership experience is itself the "
            "edge, and groups whose founder pipeline arrives through warm member "
            "referrals rather than cold applications. Brydge Club's cohort framing is "
            "the clearest example: scarcity is the message."
        ),
    )

    archetype_slide(
        prs,
        num="03",
        name="The platform or multi-product homepage",
        count="~5 sites",
        summary="Less an angel network, more a content, tools, and community brand.",
        examples="Hustle Fund / Angel Squad, The Helm, iFundWomen, She Angel Investors, Coralus.",
        why_present=(
            "When a group publishes newsletters, books, events, and editorial alongside "
            "the fund itself, the homepage has to route attention across products. The "
            "investing arm becomes one tile among many. The benefit is reach: free "
            "resources pull a wide audience and the syndicate or fund converts the "
            "narrow slice that fits."
        ),
        why_chosen=(
            "Operators with a marketing-and-media mindset, often venture studios, "
            "publishers, or hybrid fund-plus-syndicate vehicles. The Helm pairs a fund "
            "with The Helm Review. Hustle Fund pairs a pre-seed fund with Angel Squad, "
            "books, and a scrappy founder blog. Editorial volume is the trust signal."
        ),
    )

    archetype_slide(
        prs,
        num="04",
        name="The personal or portfolio site",
        count="2 sites",
        summary="A lower-key, individual-investor presence. The person is the brand.",
        examples="Gotham Gal Ventures (Joanne Wilson), The Council Angels (operator-led).",
        why_present=(
            "When the investor is the institution, the homepage is essentially a "
            "writer's website with portfolio links. There is no membership tier, no "
            "founder application form, and often no fund vehicle visible at all. The "
            "page exists so that warm intros land somewhere credible."
        ),
        why_chosen=(
            "Solo angels and very small operator-led syndicates whose deal flow runs "
            "entirely through trust networks. The footprint is intentionally small. "
            "Anyone arriving cold is meant to read, follow, and reach back through a "
            "shared connection rather than apply through a form."
        ),
    )

    slide_takeaway(prs)
    slide_standouts(prs)
    slide_close(prs)

    out_path = (
        "/Users/tori/Projects/Silver-and-Salt/_research/decks/"
        "angel-syndicate-archetypes.pptx"
    )
    prs.save(out_path)
    print(f"Wrote: {out_path}")


if __name__ == "__main__":
    main()
